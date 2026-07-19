"""Generate Valaiyagam client PPT and project tracking Excel."""

from datetime import date, timedelta
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
START = date(2026, 7, 14)


def week_date(week_offset: int, weekday: int = 0) -> date:
    """Return Monday (or weekday) of project week N (1-based)."""
    base_monday = START - timedelta(days=START.weekday())
    return base_monday + timedelta(weeks=week_offset - 1, days=weekday)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    shape.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x14, 0xB8, 0xA6)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(0x47, 0x55, 0x69)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], footer: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.55), Inches(0.12), Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x14, 0xB8, 0xA6)
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    body = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        p.space_after = Pt(10)

    if footer:
        fb = slide.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.35))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Valaiyagam E-Commerce Platform",
        "Architecture · UI Standards · Payments · Courier · Tracking · Timeline\nClient Presentation · July 2026",
    )

    add_bullets_slide(
        prs,
        "Agenda",
        [
            "1. Product vision and business outcomes",
            "2. End-to-end architecture (Admin + Store + API + MySQL)",
            "3. UI / UX standards (light glass theme)",
            "4. Payment integration workflow",
            "5. Courier integration workflow",
            "6. Order tracking system",
            "7. Delivery timeline and milestones",
            "8. Git / story tracking process",
            "9. Next steps and client inputs needed",
        ],
    )

    add_bullets_slide(
        prs,
        "Product Vision",
        [
            "One platform for storefront shopping and admin operations",
            "Secure role-based admin for catalog, orders, users, and fulfillment",
            "Trusted checkout with payment gateway + COD options",
            "Reliable courier booking, labels, and live shipment tracking",
            "Mobile-first customer and admin experiences",
            "Dockerized, migration-ready foundation for long-term growth",
        ],
        "Foundation already live: Auth, Users, Roles, Docker, Alembic, Glass Admin UI",
    )

    add_bullets_slide(
        prs,
        "Architecture Overview",
        [
            "Frontend: Next.js + React + Tailwind + Redux Toolkit",
            "Backend: FastAPI layered design (Route → Service → Repository → Model)",
            "Database: MySQL 8 with Alembic versioned migrations",
            "Infrastructure: Docker Compose network (frontend, backend, mysql)",
            "Integrations: Payment gateway APIs + Courier partner adapters",
            "Security: JWT auth, Argon2 passwords, RBAC, webhook signatures",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Visual System",
        [
            "Light theme only — soft #f5f8fc backgrounds (no dark shells)",
            "Glassmorphism panels: translucent white, blur, soft borders",
            "Primary accent: teal → cyan gradient buttons and active states",
            "Typography: clear hierarchy, generous spacing, readable contrast",
            "Icons: simple line icons for navigation and actions",
            "Motion: subtle hover lift / shadow — never noisy animation",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Layout & Interaction",
        [
            "Desktop: glass sidebar + content workspace",
            "Mobile: sticky bottom navigation + card lists (not dense tables)",
            "Create / Edit always open in glass modal forms",
            "Delete always requires a confirmation modal (never browser alerts)",
            "One primary job per screen/section",
            "Inline error banners from API; empty states with clear CTA",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Responsive Rules",
        [
            "Breakpoints cover phone, tablet, and desktop admin use",
            "Touch targets sized for thumbs on mobile actions",
            "Modals become bottom sheets on small screens",
            "Tables transform into stacked cards under tablet width",
            "Esc key closes dialogs; focus remains usable",
            "Consistent glass components reused across all modules",
        ],
    )

    add_bullets_slide(
        prs,
        "Admin Modules Roadmap",
        [
            "Done: Authentication, Users, Roles",
            "Next: Catalog (categories, products, media), Inventory",
            "Then: Orders, Addresses, Cart/Checkout",
            "Integrations: Payments, Couriers, Tracking timeline",
            "Growth: Notifications, Reports, CMS, Store settings",
            "Each module follows the same UI + API layering standards",
        ],
    )

    add_bullets_slide(
        prs,
        "Payment Integration",
        [
            "Customer checkout creates order in PENDING_PAYMENT state",
            "Backend creates gateway order/intent (Razorpay primary; Stripe optional)",
            "Customer pays via hosted checkout / SDK — no card storage in our DB",
            "Gateway webhook verifies signature → marks order PAID",
            "Failure/timeout → PAYMENT_FAILED with safe retry",
            "Admin can issue refunds; all events audited in payment_events",
            "COD supported with courier confirmation before settlement",
        ],
    )

    add_bullets_slide(
        prs,
        "Payment — Status Lifecycle",
        [
            "PENDING_PAYMENT → PAID → (eligible for fulfillment)",
            "PENDING_PAYMENT → PAYMENT_FAILED → retry / cancel",
            "PAID → REFUND_REQUESTED → PARTIALLY_REFUNDED / REFUNDED",
            "Webhook idempotency prevents double capture",
            "Nightly reconciliation job catches missed webhooks",
            "Settlement reports available in admin finance views",
        ],
    )

    add_bullets_slide(
        prs,
        "Courier Integration",
        [
            "Adapter pattern supports multiple partners under one service",
            "Candidates: Delhivery, Shiprocket, BlueDart + manual fallback",
            "After PAID (or approved COD): create shipment → get AWB",
            "Generate shipping label and schedule pickup",
            "Serviceability check by pincode/weight before booking",
            "Cancel / reassign shipment with audit trail",
            "Credentials stored as secure courier_accounts configuration",
        ],
    )

    add_bullets_slide(
        prs,
        "Courier — Operational Flow",
        [
            "1. Validate address, package weight, and serviceability",
            "2. Apply routing rules (cost, SLA, zone)",
            "3. Call partner create_shipment API",
            "4. Persist AWB + label URL on shipment record",
            "5. Print label / handoff to warehouse packing",
            "6. Sync status via partner webhook or scheduled poll",
            "7. Handle exceptions: RTO, address issue, lost package",
        ],
    )

    add_bullets_slide(
        prs,
        "Tracking System",
        [
            "Unified timeline for customer and admin",
            "Stages: Placed → Paid → Packed → Shipped → Out for delivery → Delivered",
            "Customer track page: Order ID + verified phone/email (or login)",
            "Admin order detail shows payment events + shipment checkpoints",
            "Notifications on key milestones (email/SMS/WhatsApp)",
            "Exception queue for delayed / failed deliveries",
            "Manual override allowed for ops with mandatory reason log",
        ],
    )

    add_bullets_slide(
        prs,
        "End-to-End Order Journey",
        [
            "Browse catalog → Add to cart → Checkout address",
            "Choose payment method → Pay / COD confirm",
            "Warehouse packs order → Courier booked → AWB issued",
            "Customer receives tracking link + live timeline updates",
            "Delivery / RTO / refund handled with status sync",
            "Reports update sales, fulfillment SLA, and settlements",
        ],
    )

    add_bullets_slide(
        prs,
        "16-Week Delivery Timeline",
        [
            "W1–W2  P0 Foundation — Auth, Users, Roles, Docker (DONE)",
            "W3–W5  P1 Catalog — Products, categories, inventory basics",
            "W6–W8  P2 Orders — Cart, checkout, order lifecycle",
            "W9–W10 P3 Payments — Gateway, webhooks, refunds, COD",
            "W11–W12 P4 Courier — Adapters, AWB, labels, pickup",
            "W13–W14 P5 Tracking & Notify — Timeline + alerts",
            "W15–W16 P6 Harden & Launch — UAT, performance, go-live",
        ],
    )

    add_bullets_slide(
        prs,
        "Milestones for Client Sign-off",
        [
            "M1 Foundation demo — complete",
            "M2 Catalog demo — publish products",
            "M3 Checkout demo — place unpaid order",
            "M4 Payment demo — paid order path",
            "M5 Courier demo — AWB + label from paid order",
            "M6 Tracking demo — customer timeline live",
            "M7 Go-live readiness checklist approved",
        ],
    )

    add_bullets_slide(
        prs,
        "Git & Story Tracking Process",
        [
            "Branches: main · develop · feat/* · fix/* · hotfix/*",
            "Each story maps to a feature branch and PR",
            "Excel tracker stores: ID, Date, Phase, Story, Branch, Status, Owner",
            "Definition of Done: code + review + migration (if any) + UAT note",
            "Weekly status pulled directly from PROJECT_TRACKING.xlsx",
            "No direct commits to main — PR + review required",
        ],
    )

    add_bullets_slide(
        prs,
        "Client Inputs Needed",
        [
            "Brand logo, colors confirmation, and product sample data",
            "Payment merchant account (Razorpay/Stripe) keys for staging",
            "Courier partner accounts and warehouse pickup addresses",
            "Tax rules, shipping zones, and COD eligibility policy",
            "Notification channels (email domain / SMS provider)",
            "UAT users and preferred weekly review slot",
        ],
    )

    add_bullets_slide(
        prs,
        "Summary & Ask",
        [
            "Solid technical foundation is already running in Docker",
            "UI standard is fixed: light glass, modal edit, confirm deletes, mobile-first",
            "Payment, courier, and tracking are designed as first-class modules",
            "16-week plan delivers a launchable MVP with clear milestones",
            "Ask: approve timeline, confirm gateway + courier partners, start P1 Catalog",
        ],
        "Thank you — Valaiyagam Commerce Team",
    )

    prs.save(path)


def style_header(ws, row: int, cols: int) -> None:
    fill = PatternFill("solid", fgColor="0F766E")
    font = Font(color="FFFFFF", bold=True, size=11)
    thin = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )
    for col in range(1, cols + 1):
        cell = ws.cell(row=row, column=col)
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thin


def autosize(ws, widths: dict[int, int]) -> None:
    for col, width in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width


def status_fill(status: str) -> PatternFill | None:
    mapping = {
        "Done": "D1FAE5",
        "In Progress": "FEF3C7",
        "Planned": "E0F2FE",
        "Blocked": "FEE2E2",
    }
    color = mapping.get(status)
    return PatternFill("solid", fgColor=color) if color else None


def build_excel(path: Path) -> None:
    wb = Workbook()

    # --- Stories sheet ---
    stories = wb.active
    stories.title = "Stories"

    headers = [
        "ID",
        "Date",
        "Phase",
        "Week",
        "Module",
        "Story Title",
        "User Story",
        "Git Feature / Branch",
        "Branch Type",
        "Priority",
        "Status",
        "Owner",
        "Estimate (Days)",
        "Depends On",
        "Acceptance Criteria",
        "Notes",
    ]
    stories.append(headers)
    style_header(stories, 1, len(headers))

    rows = [
        ("VL-001", week_date(1), "P0 Foundation", 1, "Platform", "Project bootstrap & Docker network",
         "As a developer, I can run the full stack with docker compose so local setup is consistent.",
         "feat/platform-docker-compose", "feat", "P0", "Done", "Backend", 2, "-",
         "Compose brings up mysql, backend, frontend on shared network", "valaiyagam-network live"),
        ("VL-002", week_date(1), "P0 Foundation", 1, "Backend", "Layered FastAPI structure",
         "As a developer, I have model/schema/repository/service/route layers for maintainable APIs.",
         "feat/backend-layered-structure", "feat", "P0", "Done", "Backend", 2, "VL-001",
         "Folders and imports follow agreed architecture", "Aligned with client architecture doc"),
        ("VL-003", week_date(1), "P0 Foundation", 1, "Auth", "JWT login + Argon2 passwords",
         "As an admin, I can securely sign in and receive an access token.",
         "feat/auth-jwt-login", "feat", "P0", "Done", "Backend", 2, "VL-002",
         "Login returns token; invalid credentials rejected", "OAuth2 password form"),
        ("VL-004", week_date(1), "P0 Foundation", 1, "RBAC", "Default roles seed (admin/manager/viewer)",
         "As the system, I seed default roles and an initial admin user on first boot.",
         "feat/rbac-default-roles-seed", "feat", "P0", "Done", "Backend", 1, "VL-003",
         "admin/manager/viewer exist; initial admin can login", "Credentials via env"),
        ("VL-005", week_date(2), "P0 Foundation", 2, "Users", "User CRUD APIs",
         "As an admin, I can create, list, update, and delete users with role assignment.",
         "feat/users-crud-api", "feat", "P0", "Done", "Backend", 2, "VL-004",
         "All CRUD endpoints work with RBAC guards", "Admin-only write"),
        ("VL-006", week_date(2), "P0 Foundation", 2, "Roles", "Role CRUD APIs",
         "As an admin, I can manage custom roles while system roles stay protected.",
         "feat/roles-crud-api", "feat", "P0", "Done", "Backend", 1, "VL-004",
         "Cannot delete/rename protected roles unsafely", "Name update supported"),
        ("VL-007", week_date(2), "P0 Foundation", 2, "Frontend", "Light glass admin shell + Redux store",
         "As an admin, I use a light glassmorphism UI with Redux-managed auth/users/roles.",
         "feat/admin-glass-ui-redux", "feat", "P0", "Done", "Frontend", 3, "VL-005",
         "Login + dashboard match UI standards", "No dark theme"),
        ("VL-008", week_date(2), "P0 Foundation", 2, "Frontend", "User/Role edit modals + confirm delete",
         "As an admin, I can edit users/roles in modals and must confirm deletes.",
         "feat/admin-edit-confirm-modals", "feat", "P0", "Done", "Frontend", 2, "VL-007",
         "Edit works; delete uses ConfirmDialog", "Mobile responsive cards"),
        ("VL-009", week_date(2), "P0 Foundation", 2, "Database", "MySQL schema + Alembic initial migration",
         "As a developer, schema changes are versioned with Alembic and SQL init scripts.",
         "feat/alembic-initial-schema", "feat", "P0", "Done", "Backend", 2, "VL-002",
         "alembic upgrade head creates users/roles/user_roles", "002_schema.sql added"),
        ("VL-010", week_date(3), "P1 Catalog", 3, "Catalog", "Category CRUD",
         "As an admin, I can create and organize product categories.",
         "feat/catalog-categories", "feat", "P1", "Planned", "Backend", 3, "VL-009",
         "Nested or flat categories with active flag", "Includes admin UI"),
        ("VL-011", week_date(3), "P1 Catalog", 3, "Catalog", "Product create/edit with media",
         "As an admin, I can publish products with images, price, and description.",
         "feat/catalog-products-media", "feat", "P1", "Planned", "Fullstack", 5, "VL-010",
         "Product visible in admin list; media upload works", "S3/local storage TBD"),
        ("VL-012", week_date(4), "P1 Catalog", 4, "Catalog", "Product variants (size/color)",
         "As an admin, I can define variants with SKU and price overrides.",
         "feat/catalog-product-variants", "feat", "P1", "Planned", "Backend", 4, "VL-011",
         "Variants persist and appear in product detail", "-"),
        ("VL-013", week_date(4), "P1 Catalog", 4, "Inventory", "Stock levels & low-stock alerts",
         "As a manager, I can monitor stock and get low-stock warnings.",
         "feat/inventory-stock-alerts", "feat", "P1", "Planned", "Fullstack", 4, "VL-012",
         "Stock updates on admin; alert threshold configurable", "-"),
        ("VL-014", week_date(5), "P1 Catalog", 5, "Storefront", "Public catalog browse & product detail",
         "As a customer, I can browse products and open product details.",
         "feat/storefront-catalog-browse", "feat", "P1", "Planned", "Frontend", 5, "VL-011",
         "Listing + detail pages responsive", "SEO-friendly routes"),
        ("VL-015", week_date(6), "P2 Orders", 6, "Cart", "Cart add/update/remove",
         "As a customer, I can manage cart items before checkout.",
         "feat/cart-management", "feat", "P1", "Planned", "Fullstack", 4, "VL-014",
         "Guest + logged-in cart supported", "-"),
        ("VL-016", week_date(6), "P2 Orders", 6, "Checkout", "Address book & shipping quote",
         "As a customer, I can select delivery address and see shipping estimate.",
         "feat/checkout-address-shipping", "feat", "P1", "Planned", "Fullstack", 5, "VL-015",
         "Address validation + quote returned", "-"),
        ("VL-017", week_date(7), "P2 Orders", 7, "Orders", "Create order lifecycle APIs",
         "As the system, I create orders with line items and status machine.",
         "feat/orders-lifecycle-api", "feat", "P0", "Planned", "Backend", 5, "VL-016",
         "Statuses: draft/pending/paid/cancelled supported", "-"),
        ("VL-018", week_date(8), "P2 Orders", 8, "Admin", "Admin order list & detail",
         "As an admin, I can view and filter orders with line-item detail.",
         "feat/admin-orders-console", "feat", "P1", "Planned", "Frontend", 4, "VL-017",
         "Filters by status/date; detail shows totals", "Glass UI pattern"),
        ("VL-019", week_date(9), "P3 Payments", 9, "Payments", "Razorpay order + checkout SDK",
         "As a customer, I can pay for an order via Razorpay checkout.",
         "feat/payments-razorpay-checkout", "feat", "P0", "Planned", "Fullstack", 5, "VL-017",
         "Sandbox payment succeeds and maps to order", "No card data stored"),
        ("VL-020", week_date(9), "P3 Payments", 9, "Payments", "Payment webhook + signature verify",
         "As the system, I verify gateway webhooks and update payment status idempotently.",
         "feat/payments-razorpay-webhook", "feat", "P0", "Planned", "Backend", 4, "VL-019",
         "Invalid signature rejected; duplicate events ignored", "payment_events audit"),
        ("VL-021", week_date(10), "P3 Payments", 10, "Payments", "Refunds & COD flow",
         "As an admin, I can refund paid orders; customers can choose COD where allowed.",
         "feat/payments-refunds-cod", "feat", "P1", "Planned", "Fullstack", 4, "VL-020",
         "Partial/full refund + COD order path verified", "-"),
        ("VL-022", week_date(11), "P4 Courier", 11, "Courier", "Courier adapter interface",
         "As a developer, I can plug multiple courier partners behind one service contract.",
         "feat/courier-adapter-interface", "feat", "P0", "Planned", "Backend", 3, "VL-021",
         "create/cancel/track/label methods defined", "Manual adapter included"),
        ("VL-023", week_date(11), "P4 Courier", 11, "Courier", "Shiprocket/Delhivery shipment create",
         "As an ops user, I can create a shipment and receive an AWB for a paid order.",
         "feat/courier-create-shipment-awb", "feat", "P0", "Planned", "Backend", 5, "VL-022",
         "AWB saved on shipment; label URL available", "Needs partner credentials"),
        ("VL-024", week_date(12), "P4 Courier", 12, "Courier", "Label print & pickup scheduling",
         "As warehouse staff, I can print labels and schedule pickups.",
         "feat/courier-labels-pickup", "feat", "P1", "Planned", "Fullstack", 4, "VL-023",
         "Label downloadable; pickup confirmation stored", "-"),
        ("VL-025", week_date(13), "P5 Tracking", 13, "Tracking", "Shipment event ingestion",
         "As the system, I ingest courier webhooks/polls into shipment_events.",
         "feat/tracking-shipment-events", "feat", "P0", "Planned", "Backend", 4, "VL-023",
         "Timeline events ordered by timestamp", "-"),
        ("VL-026", week_date(13), "P5 Tracking", 13, "Tracking", "Customer order tracking page",
         "As a customer, I can track my order with a clear status timeline.",
         "feat/tracking-customer-timeline-ui", "feat", "P0", "Planned", "Frontend", 4, "VL-025",
         "Shows Placed→Delivered stages with AWB", "OTP or login gate"),
        ("VL-027", week_date(14), "P5 Tracking", 14, "Notify", "Order/shipment notifications",
         "As a customer, I receive email/SMS updates for payment and shipment milestones.",
         "feat/notifications-order-shipment", "feat", "P1", "Planned", "Backend", 4, "VL-026",
         "Templates for paid/shipped/delivered", "Provider TBD"),
        ("VL-028", week_date(14), "P5 Tracking", 14, "Admin", "Exception queue & manual override",
         "As ops, I can manage delayed/RTO shipments and log manual overrides.",
         "feat/tracking-exception-queue", "feat", "P1", "Planned", "Fullstack", 3, "VL-025",
         "Override requires reason; audited", "-"),
        ("VL-029", week_date(15), "P6 Launch", 15, "Reports", "Sales & fulfillment dashboard",
         "As a business owner, I can view sales, payment, and delivery KPIs.",
         "feat/reports-sales-fulfillment", "feat", "P1", "Planned", "Fullstack", 4, "VL-021",
         "Daily/weekly charts and CSV export", "-"),
        ("VL-030", week_date(15), "P6 Launch", 15, "QA", "UAT script & bugfix sprint",
         "As the client, I can execute UAT and log defects against stories.",
         "fix/uat-bugfix-pass", "fix", "P0", "Planned", "QA", 5, "VL-027",
         "Critical bugs closed; UAT sign-off recorded", "-"),
        ("VL-031", week_date(16), "P6 Launch", 16, "Ops", "Staging → production cutover",
         "As the team, we deploy production with migrations, secrets, and monitoring.",
         "feat/prod-cutover-checklist", "feat", "P0", "Planned", "DevOps", 3, "VL-030",
         "HTTPS, backups, smoke tests green", "-"),
        ("VL-032", week_date(16), "P6 Launch", 16, "Docs", "Client runbook & training",
         "As the client ops team, I have a runbook for daily admin tasks.",
         "docs/client-ops-runbook", "docs", "P2", "Planned", "PM", 2, "VL-031",
         "Runbook covers users, orders, refunds, shipments", "-"),
    ]

    thin = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )

    for row in rows:
        stories.append(list(row))
        r = stories.max_row
        for c in range(1, len(headers) + 1):
            cell = stories.cell(row=r, column=c)
            cell.border = thin
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if c == 2:
                cell.number_format = "YYYY-MM-DD"
        status_cell = stories.cell(row=r, column=11)
        fill = status_fill(str(status_cell.value))
        if fill:
            status_cell.fill = fill

    autosize(
        stories,
        {
            1: 10,
            2: 12,
            3: 14,
            4: 8,
            5: 12,
            6: 36,
            7: 55,
            8: 34,
            9: 12,
            10: 10,
            11: 12,
            12: 12,
            13: 14,
            14: 12,
            15: 42,
            16: 28,
        },
    )
    stories.freeze_panes = "A2"
    stories.auto_filter.ref = f"A1:P{stories.max_row}"
    stories.row_dimensions[1].height = 30

    # --- Timeline sheet ---
    timeline = wb.create_sheet("Timeline")
    timeline_headers = [
        "Phase",
        "Week Start",
        "Week End",
        "Focus",
        "Milestone",
        "Exit Criteria",
        "Status",
    ]
    timeline.append(timeline_headers)
    style_header(timeline, 1, len(timeline_headers))
    phases = [
        ("P0 Foundation", week_date(1), week_date(2, 6), "Auth, RBAC, Docker, UI shell", "M1 Foundation demo", "Admin login + user/role CRUD", "Done"),
        ("P1 Catalog", week_date(3), week_date(5, 6), "Categories, products, inventory", "M2 Catalog demo", "Publish products to storefront", "Planned"),
        ("P2 Orders", week_date(6), week_date(8, 6), "Cart, checkout, order lifecycle", "M3 Checkout demo", "Place order end-to-end", "Planned"),
        ("P3 Payments", week_date(9), week_date(10, 6), "Gateway, webhooks, refunds, COD", "M4 Payment demo", "Paid order verified", "Planned"),
        ("P4 Courier", week_date(11), week_date(12, 6), "Adapters, AWB, labels, pickup", "M5 Courier demo", "Shipment created from paid order", "Planned"),
        ("P5 Tracking", week_date(13), week_date(14, 6), "Timeline UI + notifications", "M6 Tracking demo", "Customer can track order", "Planned"),
        ("P6 Launch", week_date(15), week_date(16, 6), "Reports, UAT, production cutover", "M7 Go-live", "Production ready checklist", "Planned"),
    ]
    for phase in phases:
        timeline.append(list(phase))
        r = timeline.max_row
        for c in (2, 3):
            timeline.cell(row=r, column=c).number_format = "YYYY-MM-DD"
        fill = status_fill(str(phase[-1]))
        if fill:
            timeline.cell(row=r, column=7).fill = fill
    autosize(timeline, {1: 14, 2: 12, 3: 12, 4: 36, 5: 22, 6: 36, 7: 12})
    timeline.freeze_panes = "A2"

    # --- Branches sheet ---
    branches = wb.create_sheet("Git Branches")
    branch_headers = ["Branch Name", "Type", "Linked Story ID", "Title", "Target", "Status"]
    branches.append(branch_headers)
    style_header(branches, 1, len(branch_headers))
    for row in rows:
        story_id, _, _, _, _, title, _, branch, btype, _, status, *_ = row
        branches.append([branch, btype, story_id, title, "develop", status])
        fill = status_fill(status)
        if fill:
            branches.cell(row=branches.max_row, column=6).fill = fill
    autosize(branches, {1: 38, 2: 10, 3: 14, 4: 42, 5: 12, 6: 12})
    branches.freeze_panes = "A2"
    branches.auto_filter.ref = f"A1:F{branches.max_row}"

    # --- Workflow sheet ---
    workflow = wb.create_sheet("Workflow")
    workflow_headers = ["Step", "Stage", "Owner", "Input", "Output", "Tool / Gate"]
    workflow.append(workflow_headers)
    style_header(workflow, 1, len(workflow_headers))
    workflow_rows = [
        (1, "Story drafted", "PM", "Client requirement", "Story in Stories sheet", "PROJECT_TRACKING.xlsx"),
        (2, "Branch created", "Dev", "Story ID + title", "feat/fix branch", "Git"),
        (3, "Implementation", "Dev", "Acceptance criteria", "Code + tests", "IDE / Docker"),
        (4, "Migration (if schema)", "Backend", "Model change", "Alembic revision", "alembic revision"),
        (5, "Pull request", "Dev", "Branch commits", "PR to develop", "GitHub/GitLab"),
        (6, "Code review", "Peer", "PR", "Approved PR", "Review checklist"),
        (7, "CI / smoke", "CI", "PR build", "Green checks", "Docker compose tests"),
        (8, "Merge", "Maintainer", "Approved PR", "Updated develop", "Protected branch"),
        (9, "UAT on staging", "Client + QA", "Build on staging", "Pass/fail notes", "Staging URL"),
        (10, "Story status update", "PM", "UAT result", "Status=Done", "Excel tracker"),
        (11, "Release to main", "DevOps", "Stable develop", "Production tag", "Release checklist"),
    ]
    for item in workflow_rows:
        workflow.append(list(item))
    autosize(workflow, {1: 8, 2: 22, 3: 14, 4: 22, 5: 22, 6: 24})

    # --- Integrations sheet ---
    integrations = wb.create_sheet("Integrations")
    int_headers = ["Integration", "Provider Options", "Phase", "Story IDs", "Key Objects", "Success Metric"]
    integrations.append(int_headers)
    style_header(integrations, 1, len(int_headers))
    integrations.append(
        [
            "Payments",
            "Razorpay (primary), Stripe (optional), COD",
            "P3",
            "VL-019, VL-020, VL-021",
            "payments, payment_events, refunds",
            "Paid order confirmed via verified webhook",
        ]
    )
    integrations.append(
        [
            "Courier",
            "Shiprocket / Delhivery / BlueDart + manual",
            "P4",
            "VL-022, VL-023, VL-024",
            "shipments, courier_accounts, labels",
            "AWB created from paid order < 2 minutes",
        ]
    )
    integrations.append(
        [
            "Tracking",
            "Partner webhooks + internal timeline",
            "P5",
            "VL-025, VL-026, VL-027, VL-028",
            "shipment_events, notifications",
            "Customer sees live timeline with AWB",
        ]
    )
    autosize(integrations, {1: 14, 2: 42, 3: 10, 4: 28, 5: 40, 6: 42})

    # --- Dashboard sheet ---
    dash = wb.create_sheet("Dashboard", 0)
    dash["A1"] = "Valaiyagam E-Commerce — Project Tracking Dashboard"
    dash["A1"].font = Font(size=16, bold=True, color="0F766E")
    dash.merge_cells("A1:F1")

    dash["A3"] = "Project"
    dash["B3"] = "Valaiyagam Commerce Admin + Storefront"
    dash["A4"] = "Start Date"
    dash["B4"] = START
    dash["B4"].number_format = "YYYY-MM-DD"
    dash["A5"] = "Planned End"
    dash["B5"] = week_date(16, 6)
    dash["B5"].number_format = "YYYY-MM-DD"
    dash["A6"] = "Document Date"
    dash["B6"] = date(2026, 7, 19)
    dash["B6"].number_format = "YYYY-MM-DD"

    done = sum(1 for r in rows if r[10] == "Done")
    planned = sum(1 for r in rows if r[10] == "Planned")
    total = len(rows)

    dash["A8"] = "KPI"
    dash["B8"] = "Value"
    style_header(dash, 8, 2)
    dash["A9"] = "Total Stories"
    dash["B9"] = total
    dash["A10"] = "Done"
    dash["B10"] = done
    dash["A11"] = "Planned"
    dash["B11"] = planned
    dash["A12"] = "% Complete"
    dash["B12"] = done / total
    dash["B12"].number_format = "0.0%"

    dash["A14"] = "How to use this workbook"
    dash["A14"].font = Font(bold=True, size=12)
    instructions = [
        "1. Update Stories sheet status weekly (Done / In Progress / Planned / Blocked).",
        "2. Create Git branch using the exact name in column Git Feature / Branch.",
        "3. Link PRs to Story ID (example: VL-019) in the PR title.",
        "4. Review Timeline sheet in client weekly calls.",
        "5. Use Workflow sheet as Definition of Done checklist.",
        "6. Track payment/courier/tracking readiness on Integrations sheet.",
    ]
    for i, text in enumerate(instructions):
        dash.cell(row=15 + i, column=1, value=text)

    dash["A22"] = "Related docs"
    dash["A22"].font = Font(bold=True)
    dash["A23"] = "docs/ECOM_ARCHITECTURE_AND_TIMELINE.md"
    dash["A24"] = "docs/Valaiyagam_Client_Presentation.pptx"

    autosize(dash, {1: 55, 2: 45})

    wb.save(path)


def main() -> None:
    pptx_path = ROOT / "Valaiyagam_Client_Presentation.pptx"
    xlsx_path = ROOT / "PROJECT_TRACKING.xlsx"
    build_pptx(pptx_path)
    build_excel(xlsx_path)
    print(f"Wrote {pptx_path}")
    print(f"Wrote {xlsx_path}")


if __name__ == "__main__":
    main()
